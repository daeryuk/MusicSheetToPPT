from flask import Flask, render_template, request, jsonify, send_file
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import re
import webbrowser
from threading import Timer
from io import BytesIO

app = Flask(__name__)

def clean_lyrics(lyrics):
    # _x000D_ 를 실제 줄바꿈으로 변환
    cleaned = lyrics.replace('_x000D_', '\n')
    # 불필요한 문자 및 패턴 제거
    cleaned = re.sub(r'^[•●■□○\-\*]\s*', '', cleaned, flags=re.MULTILINE)  # 글머리 기호 제거
    cleaned = re.sub(r'\s*×\d+\s*', '', cleaned)  # ×2, ×3 등의 반복 표시 제거
    cleaned = re.sub(r'[.]{3,}', '', cleaned)  # 점(...) 제거
    
    # 3개 이상의 연속된 줄바꿈을 2개로 통일
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    
    # 빈 줄만 있는 라인 제거
    lines = [line for line in cleaned.split('\n') if line.strip()]
    cleaned = '\n'.join(lines)
    
    # 시작과 끝의 불필요한 공백 제거
    cleaned = cleaned.strip()
    
    return cleaned

def search_lyrics(song_title):
    search_url = f"https://music.bugs.co.kr/search/track?q={song_title}"
    
    try:
        response = requests.get(search_url)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        track_info = soup.find('a', {'class': 'trackInfo'})
        if not track_info:
            return None, None
            
        song_url = track_info['href']
        song_page = requests.get(song_url)
        song_soup = BeautifulSoup(song_page.text, 'html.parser')
        
        lyrics_container = song_soup.find('div', {'class': 'lyricsContainer'})
        if lyrics_container and lyrics_container.find('xmp'):
            lyrics = lyrics_container.find('xmp').text.strip()
            # 가사 전처리 적용
            lyrics = clean_lyrics(lyrics)
            return lyrics, song_url
        return None, None
    except:
        return None, None

def create_ppt(song_order, memory_file=True):
    """
    PPT 생성 함수
    Args:
        song_order: 곡 정보가 담긴 리스트
        memory_file: True면 BytesIO 객체 반환, False면 파일로 저장
    """
    prs = Presentation()

    for song_idx in song_order:
        song_title = song_idx['title']
        lyrics = song_idx['lyrics']
        
        # 제목 슬라이드 생성
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = song_title
        
        # 제목 스타일 설정
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(44)

        # 문단을 분리하여 슬라이드 생성
        paragraphs = lyrics.split('\n\n')  # 빈 줄을 기준으로 문단 분리
        
        for paragraph in paragraphs:
            if paragraph.strip():  # 빈 문단이 아닌 경우에만 슬라이드 생성
                content_slide = prs.slides.add_slide(prs.slide_layouts[5])
                text_box = content_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
                text_frame = text_box.text_frame
                text_frame.text = paragraph.strip()

                # 텍스트 스타일 설정
                for para in text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    para.font.size = Pt(32)

    # 메모리에 PPT 파일을 저장하고 반환
    pptx_buffer = BytesIO()  # 메모리 버퍼 생성
    prs.save(pptx_buffer)    # PPT를 메모리에 저장
    pptx_buffer.seek(0)      # 버퍼의 포인터를 처음으로 이동
    return pptx_buffer       # 메모리에 저장된 PPT 반환


@app.route('/create_ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.json
        pptx_buffer = create_ppt(data['songOrder'], memory_file=True)
        
        return send_file(
            pptx_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='worship_lyrics.pptx'
        )
    except Exception as e:
        print(f"PPT 생성 중 에러 발생: {str(e)}")
        return jsonify({'error': str(e)}), 500
        
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/search_lyrics', methods=['POST'])
def process_lyrics():
    songs = request.json['songs']
    results = {}
    
    for song in songs:
        lyrics, song_url = search_lyrics(song)
        if lyrics:
            results[song] = {
                'lyrics': lyrics,
                'url': song_url
            }
        else:
            results[song] = {
                'lyrics': "가사를 찾을 수 없습니다.",
                'url': None
            }
    
    return jsonify(results)

if __name__ == '__main__':
    app.run(debug=True)    
